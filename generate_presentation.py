from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import json

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_section_slide(prs, title):
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    return slide

def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
    return slide

def add_content_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = content
    p.font.size = Pt(20)
    return slide

def add_image_slide(prs, title, image_path, caption=None):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(7)
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width=width)
    if caption:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(7), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
    return slide

def get_summary():
    with open('data_analysis_summary.txt', 'r', encoding='utf-8') as f:
        return f.read()

def get_outline():
    with open('tata_steel_presentation_outline.txt', 'r', encoding='utf-8') as f:
        return f.read()

def parse_outline_sections(outline):
    slides = []
    current_title = None
    current_content = []
    for line in outline.splitlines():
        if line.startswith('## Slide'):
            if current_title:
                slides.append((current_title, current_content))
            current_title = line.replace('## ', '').strip()
            current_content = []
        elif line.startswith('- '):
            current_content.append(line[2:].strip())
        elif line.strip() and not line.startswith('---'):
            current_content.append(line.strip())
    if current_title:
        slides.append((current_title, current_content))
    return slides

def main():
    prs = Presentation()
    # Title Slide
    add_title_slide(prs, "Tata Steel Workforce Resource Analysis & Utilization Applications", "Data-Driven Insights for Organizational Optimization")

    # Parse outline for slides
    outline = get_outline()
    slides = parse_outline_sections(outline)

    # Map slide titles to chart images
    chart_map = {
        'Division Distribution Analysis': 'charts/division_distribution.png',
        'Resource Concentration': 'charts/resource_distribution_pie.png',
        'Group Analysis': 'charts/group_distribution.png',
        'Key Division-Group Relationships': 'charts/division_group_relationships.png',
        'Department Analysis': 'charts/department_distribution.png',
    }

    for title, content in slides:
        # If the slide is a chart slide
        for key, img in chart_map.items():
            if key in title:
                add_image_slide(prs, title, img)
                break
        else:
            # If the slide is a bullet or content slide
            if len(content) > 1:
                add_bullet_slide(prs, title, content)
            elif len(content) == 1:
                add_content_slide(prs, title, content[0])
            else:
                add_section_slide(prs, title)

    # Save the presentation
    prs.save('Tata_Steel_Utilization_Applications.pptx')
    print('Presentation created: Tata_Steel_Utilization_Applications.pptx')

if __name__ == "__main__":
    main() 